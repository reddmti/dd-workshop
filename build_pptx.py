"""Genera WORKSHOP_APERTURA.pptx - presentacion de apertura para la workshop."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = r"C:\Users\matia\Documents\GitHub\dd-workshop\WORKSHOP_APERTURA.pptx"
LOGO = r"C:\Users\matia\Documents\GitHub\dd-workshop\logo.jpg"

# Colores corporativos / Datadog
PURPLE = RGBColor(0x63, 0x2C, 0xA6)
DARK = RGBColor(0x1F, 0x1F, 0x1F)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF5, 0xF0, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
BLUE = RGBColor(0x15, 0x65, 0xC0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, size=18, color=DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = "  " + item
        r.font.name = 'Calibri'
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def add_logo(slide, right=True):
    # Logo TekServicios
    logo_w = Inches(1.4)
    if right:
        slide.shapes.add_picture(LOGO, SW - logo_w - Inches(0.3),
                                 Inches(0.25), height=Inches(0.55))
    else:
        slide.shapes.add_picture(LOGO, Inches(0.3),
                                 Inches(0.25), height=Inches(0.55))

def add_footer(slide, page_num, total):
    add_rect(slide, 0, SH - Inches(0.35), SW, Inches(0.35), PURPLE)
    add_text(slide, Inches(0.4), SH - Inches(0.32), Inches(8), Inches(0.3),
             "Workshop Datadog  |  TekServicios", size=10, color=WHITE)
    add_text(slide, SW - Inches(1.5), SH - Inches(0.32), Inches(1.2), Inches(0.3),
             f"{page_num} / {total}", size=10, color=WHITE, align=PP_ALIGN.RIGHT)

TOTAL = 6

# ============ SLIDE 1: PORTADA ============
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, 0, Inches(0.3), SH, PURPLE)

# Logo grande centrado arriba
s.shapes.add_picture(LOGO, Inches(5.5), Inches(0.8), height=Inches(1.0))

add_text(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.3),
         "Calidad y velocidad", size=54, bold=True, color=PURPLE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.7), Inches(11.3), Inches(1.0),
         "en el ciclo de entrega", size=44, bold=True, color=DARK,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.0), Inches(11.3), Inches(0.6),
         "Una mirada practica con Datadog", size=22, color=GREY,
         align=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
         "TekServicios  |  Workshop tecnico", size=14, color=PURPLE,
         align=PP_ALIGN.CENTER)

# ============ SLIDE 2: EL PROBLEMA ============
s = prs.slides.add_slide(blank)
add_logo(s)
add_rect(s, 0, 0, Inches(0.3), SH, RED)
add_text(s, Inches(0.7), Inches(0.9), Inches(11), Inches(0.5),
         "EL PROBLEMA", size=14, bold=True, color=RED)
add_text(s, Inches(0.7), Inches(1.4), Inches(11), Inches(1.0),
         "Cuando el pipeline no protege", size=40, bold=True, color=DARK)

add_bullets(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(3.5), [
    "Tests rotos llegan a produccion",
    "Vulnerabilidades pasan inadvertidas durante semanas",
    "Cada incidente se detecta tarde y caro",
    "Las metricas de entrega se deterioran sin alerta",
], size=22)

add_footer(s, 2, TOTAL)

# ============ SLIDE 3: LA PROPUESTA ============
s = prs.slides.add_slide(blank)
add_logo(s)
add_rect(s, 0, 0, Inches(0.3), SH, PURPLE)
add_text(s, Inches(0.7), Inches(0.9), Inches(11), Inches(0.5),
         "LA PROPUESTA", size=14, bold=True, color=PURPLE)
add_text(s, Inches(0.7), Inches(1.4), Inches(11), Inches(1.0),
         "Dos repos, dos realidades", size=40, bold=True, color=DARK)

# Dos columnas comparativas
col_y = Inches(2.8)
col_h = Inches(3.6)
col_w = Inches(5.8)

# Columna 1 - con gates
add_rect(s, Inches(0.7), col_y, col_w, Inches(0.6), GREEN)
add_text(s, Inches(0.7), col_y, col_w, Inches(0.6),
         "REPO 1: con gates", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, Inches(0.7), col_y + Inches(0.6), col_w, col_h - Inches(0.6),
         LIGHT)
add_bullets(s, Inches(0.95), col_y + Inches(0.85), col_w - Inches(0.5),
            col_h - Inches(0.8), [
    "Tests bloquean el deploy",
    "PR Gate impide el merge",
    "El bug nunca llega a usuarios",
    "Datadog PREVIENE",
], size=16)

# Columna 2 - sin gates
add_rect(s, Inches(6.85), col_y, col_w, Inches(0.6), RED)
add_text(s, Inches(6.85), col_y, col_w, Inches(0.6),
         "REPO 2: sin gates", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, Inches(6.85), col_y + Inches(0.6), col_w, col_h - Inches(0.6),
         LIGHT)
add_bullets(s, Inches(7.1), col_y + Inches(0.85), col_w - Inches(0.5),
            col_h - Inches(0.8), [
    "Pipeline despliega siempre",
    "El bug llega a produccion",
    "Sintetico detecta la caida",
    "Datadog DETECTA",
], size=16)

add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
         "Mismo codigo, mismo bug, dos historias diferentes",
         size=14, color=GREY, align=PP_ALIGN.CENTER)

add_footer(s, 3, TOTAL)

# ============ SLIDE 4: LO QUE VAMOS A VER ============
s = prs.slides.add_slide(blank)
add_logo(s)
add_rect(s, 0, 0, Inches(0.3), SH, PURPLE)
add_text(s, Inches(0.7), Inches(0.9), Inches(11), Inches(0.5),
         "LO QUE VAMOS A VER", size=14, bold=True, color=PURPLE)
add_text(s, Inches(0.7), Inches(1.4), Inches(11), Inches(1.0),
         "Cinco capacidades en accion", size=40, bold=True, color=DARK)

# 5 cards horizontales
cards = [
    ("Code Security", "Detecta vulnerabilidades sin ejecutar"),
    ("Test Optimization", "Visibilidad real de los tests"),
    ("Continuous Testing", "Valida la app desplegada"),
    ("DORA Metrics", "Mide el impacto en la entrega"),
    ("PR Gates", "Bloquea lo que no debe llegar"),
]
card_w = Inches(2.4)
gap = Inches(0.15)
total_w = card_w * 5 + gap * 4
start_x = (SW - total_w) / 2
card_y = Inches(3.0)
card_h = Inches(3.0)

for i, (title, desc) in enumerate(cards):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, card_y, card_w, Inches(0.7), PURPLE)
    add_text(s, x, card_y, card_w, Inches(0.7),
             f"{i+1}", size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, card_y + Inches(0.7), card_w, card_h - Inches(0.7),
             LIGHT)
    add_text(s, x + Inches(0.15), card_y + Inches(0.95), card_w - Inches(0.3),
             Inches(0.8), title, size=15, bold=True, color=DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), card_y + Inches(1.7), card_w - Inches(0.3),
             Inches(1.2), desc, size=12, color=GREY,
             align=PP_ALIGN.CENTER)

add_footer(s, 4, TOTAL)

# ============ SLIDE 5: QUE SE LLEVAN ============
s = prs.slides.add_slide(blank)
add_logo(s)
add_rect(s, 0, 0, Inches(0.3), SH, PURPLE)
add_text(s, Inches(0.7), Inches(0.9), Inches(11), Inches(0.5),
         "VALOR CONCRETO", size=14, bold=True, color=PURPLE)
add_text(s, Inches(0.7), Inches(1.4), Inches(11), Inches(1.0),
         "Que se llevan al final", size=40, bold=True, color=DARK)

benefits = [
    ("Detectar antes que el usuario",
     "Cada problema lo encuentra el equipo, no el cliente"),
    ("Bloquear merges riesgosos",
     "Reglas automaticas, sin depender del criterio del revisor"),
    ("Medir el impacto de cada deploy",
     "Metricas DORA actualizadas en cada despliegue"),
    ("Defender la inversion en calidad",
     "Datos concretos para sustentar decisiones tecnicas"),
]
b_y = Inches(2.9)
b_h = Inches(0.95)
for i, (h, d) in enumerate(benefits):
    y = b_y + b_h * Inches(0).__class__(0) + Emu(int(b_h) * i + int(Emu(0)))
    y = b_y + Emu(int(b_h) * i)
    add_rect(s, Inches(0.7), y + Emu(int(Inches(0.05))),
             Inches(0.15), b_h - Inches(0.15), PURPLE)
    add_text(s, Inches(1.0), y + Inches(0.05), Inches(11.5), Inches(0.45),
             h, size=20, bold=True, color=DARK)
    add_text(s, Inches(1.0), y + Inches(0.5), Inches(11.5), Inches(0.4),
             d, size=14, color=GREY)

add_footer(s, 5, TOTAL)

# ============ SLIDE 6: AGENDA ============
s = prs.slides.add_slide(blank)
add_logo(s)
add_rect(s, 0, 0, Inches(0.3), SH, PURPLE)
add_text(s, Inches(0.7), Inches(0.9), Inches(11), Inches(0.5),
         "AGENDA", size=14, bold=True, color=PURPLE)
add_text(s, Inches(0.7), Inches(1.4), Inches(11), Inches(1.0),
         "Lo que viene en los proximos 25 minutos", size=34, bold=True,
         color=DARK)

# Tres bloques timeline
blocks = [
    ("Bloque 1", "Repo con gates", "SAST + Tests + PR Gates", "13 min", GREEN),
    ("Bloque 2", "Repo sin gates", "Continuous Testing", "5 min", RED),
    ("Bloque 3", "Comparacion", "DORA Metrics", "4 min", BLUE),
]
bk_y = Inches(3.0)
bk_h = Inches(3.0)
bk_w = Inches(3.9)
gap_b = Inches(0.25)
total_b = bk_w * 3 + gap_b * 2
sx = (SW - total_b) / 2

for i, (label, repo, demos, t, color) in enumerate(blocks):
    x = sx + (bk_w + gap_b) * i
    add_rect(s, x, bk_y, bk_w, Inches(0.65), color)
    add_text(s, x, bk_y, bk_w, Inches(0.65), label, size=16, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, bk_y + Inches(0.65), bk_w, bk_h - Inches(0.65), LIGHT)
    add_text(s, x + Inches(0.2), bk_y + Inches(0.85),
             bk_w - Inches(0.4), Inches(0.5),
             repo, size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), bk_y + Inches(1.45),
             bk_w - Inches(0.4), Inches(0.8),
             demos, size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), bk_y + Inches(2.4),
             bk_w - Inches(0.4), Inches(0.5),
             t, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
         "Vamos con la primera demo.",
         size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

add_footer(s, 6, TOTAL)

prs.save(OUT)
print(f"OK: {OUT}")
